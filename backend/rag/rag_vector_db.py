import threading
import pymupdf4llm
import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_core.documents import Document
from langchain_experimental.text_splitter import SemanticChunker
from langchain_openai import OpenAIEmbeddings
from langchain_chroma import Chroma
from dotenv import load_dotenv
import os
load_dotenv()

_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
_PERSIST_DIR = os.path.join(_ROOT, "chroma_db")
_SHARED_DIR = os.path.join(_PERSIST_DIR, "shared")

def _get_embedding_model():
    return OpenAIEmbeddings(
        model="openai/text-embedding-3-small",
        base_url="https://openrouter.ai/api/v1",
        api_key=os.getenv("OPENROUTER_API_KEY"))


def _load_pdf(file_path: str) -> str:
    """Convert PDF to Markdown via pymupdf4llm, preserving list and heading structure."""
    return pymupdf4llm.to_markdown(file_path)


def _load_docx(file_path: str) -> str:
    """Extract paragraph and table text from a Word document."""
    doc = DocxDocument(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    """Extract slide text from a PowerPoint presentation, one sentence-terminated block per slide."""
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text_frame.text for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            parts.append(f"Slide {i}: " + ". ".join(texts) + ".")
    return "\n".join(parts)


def _load_xlsx(file_path: str) -> str:
    """Extract cell text from an Excel workbook, one sentence-terminated block per row."""
    wb = openpyxl.load_workbook(file_path, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"Sheet: {sheet.title}.")
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                parts.append(" | ".join("" if cell is None else str(cell) for cell in row) + ".")
    return "\n".join(parts)


def _load_text(file_path: str) -> str:
    """Read a plain text or Markdown file as-is."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


_LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".xlsx": _load_xlsx,
    ".txt": _load_text,
    ".md": _load_text,
}

SUPPORTED_EXTENSIONS = set(_LOADERS)


def _load_document(file_path: str) -> list[Document]:
    """Dispatch to the loader for this file's extension and wrap the extracted text in a Document."""
    ext = os.path.splitext(file_path)[1].lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported file type: {ext}")
    text = loader(file_path)
    if not text.strip():
        raise ValueError(f"No text could be extracted from {os.path.basename(file_path)}")
    return [Document(page_content=text, metadata={"source": file_path})]


def delete_document(file_path: str, session_id: str, persist_directory: str) -> int:
    """Delete all ChromaDB chunks whose session_id and source metadata match. Returns deleted count."""
    vectorstore = Chroma(persist_directory=persist_directory, embedding_function=_get_embedding_model())
    results = vectorstore.get(where={"$and": [{"session_id": session_id}, {"source": file_path}]})
    ids = results["ids"]
    if ids:
        vectorstore.delete(ids)
    return len(ids)


def get_shared_persist_dir() -> str:
    return _SHARED_DIR


def add_document_for_session(file_path: str, session_id: str, user_id: str = None,
                             cancel_event: threading.Event = None) -> int:
    """Embed a document into the shared ChromaDB, tagged with session_id/document_name/user_id metadata. Returns chunk count."""
    os.makedirs(_SHARED_DIR, exist_ok=True)

    docs = _load_document(file_path)
    docs[0].metadata["session_id"] = session_id
    docs[0].metadata["document_name"] = os.path.basename(file_path)
    if user_id:
        docs[0].metadata["user_id"] = user_id
    chunks = SemanticChunker(_get_embedding_model()).split_documents(docs)

    if cancel_event is not None and cancel_event.is_set():
        return 0

    delete_document(file_path, session_id, _SHARED_DIR)  # replace any prior ingestion of this same file instead of accumulating duplicates
    vectorstore = Chroma(persist_directory=_SHARED_DIR, embedding_function=_get_embedding_model())
    vectorstore.add_documents(chunks)
    return len(chunks)


def delete_session_vectorstore(session_id: str) -> None:
    """Remove all chunks belonging to a session from the shared ChromaDB."""
    vectorstore = Chroma(persist_directory=_SHARED_DIR, embedding_function=_get_embedding_model())
    results = vectorstore.get(where={"session_id": session_id})
    ids = results["ids"]
    if ids:
        vectorstore.delete(ids)
