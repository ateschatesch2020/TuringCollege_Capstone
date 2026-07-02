import logging
import json
import os
import time
import uuid
from langchain_core.tools import tool
from langchain_core.runnables import RunnableConfig
from langchain_chroma import Chroma
from dotenv import load_dotenv
load_dotenv()
import serpapi

logger = logging.getLogger(__name__)

_GENERATED_FILES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_files")
_API_BASE_URL = "http://localhost:8001"


def _ensure_generated_dir():
    os.makedirs(_GENERATED_FILES_DIR, exist_ok=True)


def _file_size_str(path: str) -> str:
    try:
        b = os.path.getsize(path)
    except OSError:
        return "?"
    if b < 1_048_576:
        return f"{b / 1024:.1f} KB"
    return f"{b / 1_048_576:.1f} MB"


def _build_file_select_block(paths: list) -> str:
    """Returns a ```file-select JSON block for PDF files (max 10), or empty string.
    If a path is a directory, walks into it to collect PDFs."""
    collected = []
    for p in paths:
        if os.path.isdir(p):
            for root, _, files in os.walk(p):
                for f in files:
                    if f.lower().endswith(".pdf"):
                        collected.append(os.path.join(root, f))
        elif p.lower().endswith(".pdf"):
            collected.append(p)
    pdf_paths = collected[:10]
    if not pdf_paths:
        return ""
    items = []
    for p in pdf_paths:
        try:
            size = os.path.getsize(p)
        except OSError:
            size = 0
        items.append({"path": p, "name": os.path.basename(p), "size_bytes": size})
    return "\n\n```file-select\n" + json.dumps(items, ensure_ascii=False) + "\n```"


@tool
def web_search(query: str) -> str:
    """Search the web for current information using Google Search.
    Use this for any question requiring up-to-date information not available in uploaded documents.
    query: search query in any language
    Returns top search results with title, snippet, and link.
    """
    client = serpapi.Client(api_key=os.getenv("SERPAPI_KEY"))
    try:
        raw = client.search({"engine": "google", "q": query, "num": 5})
    except serpapi.exceptions.HTTPError as e:
        logger.warning("web_search: SerpAPI error for '%s': %s", query, e)
        return f"Web search unavailable: {e}"
    results = raw.get("organic_results", [])
    if not results:
        return "No results found."
    lines = []
    for r in results[:5]:
        lines.append(f"**{r.get('title', '')}**\n{r.get('snippet', '')}\n{r.get('link', '')}")
    return "\n\n".join(lines)


@tool
def generate_presentation(title: str, slides_json: str) -> str:
    """Create a PowerPoint presentation (.pptx) file and return a download link.
    Use this when the user asks to create a presentation or slideshow.
    title: presentation title
    slides_json: JSON array of slides, e.g.:
                 '[{"title": "Introduction", "content": "Point 1\\nPoint 2\\nPoint 3"}]'
                 Each slide must have a 'title' and 'content' (newline-separated bullet points).
    Returns a download URL for the generated .pptx file.
    """
    from pptx import Presentation

    try:
        slides_data = json.loads(slides_json)
    except json.JSONDecodeError as e:
        return f"Invalid slides_json: {e}. Expected JSON array like [{{\"title\": \"...\", \"content\": \"...\"}}]"

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = ""

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get("title", "")
        tf = slide.placeholders[1].text_frame
        tf.clear()
        bullet_lines = [l.strip() for l in slide_data.get("content", "").split("\n") if l.strip()]
        for i, line in enumerate(bullet_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                tf.add_paragraph().text = line

    _ensure_generated_dir()
    filename = f"{uuid.uuid4().hex}.pptx"
    prs.save(os.path.join(_GENERATED_FILES_DIR, filename))
    logger.info("generate_presentation: created %s", filename)
    return f"Presentation created successfully.\n\n[Download {title}.pptx]({_API_BASE_URL}/files/{filename})"


@tool
def generate_word_document(title: str, content: str) -> str:
    """Create a Word document (.docx) file and return a download link.
    Use this when the user asks to create a Word document, report, summary, checklist, or price list as a file.
    title: document title (used as the main heading)
    content: document body. Supports markdown-like formatting:
             - Lines starting with # / ## / ### are headings (level 1/2/3)
             - Lines starting with - or * are bullet list items
             - Lines with | delimiters are table rows (consecutive | rows become one table)
             - All other lines are normal paragraphs
    Returns a download URL for the generated .docx file.
    """
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)

    lines = content.split("\n")
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        if stripped.startswith("|"):
            table_rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                row = lines[i].strip()
                if not all(c in "-| :" for c in row):
                    table_rows.append([c.strip() for c in row.strip("|").split("|")])
                i += 1
            if table_rows:
                max_cols = max(len(r) for r in table_rows)
                tbl = doc.add_table(rows=len(table_rows), cols=max_cols)
                tbl.style = "Table Grid"
                for r_idx, cells in enumerate(table_rows):
                    for c_idx, text in enumerate(cells):
                        if c_idx < max_cols:
                            tbl.rows[r_idx].cells[c_idx].text = text
            continue

        if not stripped:
            i += 1
            continue

        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        else:
            doc.add_paragraph(stripped)
        i += 1

    _ensure_generated_dir()
    filename = f"{uuid.uuid4().hex}.docx"
    doc.save(os.path.join(_GENERATED_FILES_DIR, filename))
    logger.info("generate_word_document: created %s", filename)
    return f"Word document created successfully.\n\n[Download {title}.docx]({_API_BASE_URL}/files/{filename})"


@tool
def generate_pdf_document(title: str, content: str) -> str:
    """Create a PDF document and return a download link.
    Use this when the user asks to create a PDF report, price list, or summary as a file.
    title: document title
    content: document body. Supports:
             - Lines starting with # / ## / ### for headings
             - Lines starting with - or * for bullet points
             - Blank lines for paragraph spacing
    Returns a download URL for the generated .pdf file.
    """
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(0, 10, title, align="C")
    pdf.ln(6)
    pdf.set_font("Helvetica", "", 11)

    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            pdf.ln(3)
        elif stripped.startswith("### "):
            pdf.set_font("Helvetica", "B", 12)
            pdf.multi_cell(0, 7, stripped[4:])
            pdf.set_font("Helvetica", "", 11)
        elif stripped.startswith("## "):
            pdf.set_font("Helvetica", "B", 13)
            pdf.multi_cell(0, 8, stripped[3:])
            pdf.set_font("Helvetica", "", 11)
        elif stripped.startswith("# "):
            pdf.set_font("Helvetica", "B", 14)
            pdf.multi_cell(0, 9, stripped[2:])
            pdf.set_font("Helvetica", "", 11)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            pdf.multi_cell(0, 7, f"  - {stripped[2:]}")
        else:
            pdf.multi_cell(0, 7, stripped)

    _ensure_generated_dir()
    filename = f"{uuid.uuid4().hex}.pdf"
    pdf.output(os.path.join(_GENERATED_FILES_DIR, filename))
    logger.info("generate_pdf_document: created %s", filename)
    return f"PDF document created successfully.\n\n[Download {title}.pdf]({_API_BASE_URL}/files/{filename})"


@tool
def search_project_files(project_name: str, query: str) -> str:
    """Find a project folder by name on disk and search its file contents for a keyword.
    Use this when the user asks to find files in a project, search code or content within a project directory,
    or wants to know what files exist in a named project.
    project_name: name of the project folder (searched under the PROJECTS_DIR environment variable)
    query: keyword or text to search for inside files
    Returns the file listing and content matches with file path and line number.
    """
    projects_dir = os.getenv("PROJECTS_DIR")
    if not projects_dir:
        return "PROJECTS_DIR environment variable is not set. Cannot search project files."

    exact_match = None
    partial_matches = []
    try:
        for root, dirs, _ in os.walk(projects_dir):
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for d in list(dirs):
                if d.lower() == project_name.lower():
                    exact_match = os.path.join(root, d)
                    dirs.clear()
                    break
                if project_name.lower() in d.lower():
                    partial_matches.append(os.path.join(root, d))
            if exact_match:
                break
    except PermissionError as e:
        return f"Permission denied accessing {projects_dir}: {e}"

    if exact_match:
        project_paths = [exact_match]
    elif partial_matches:
        project_paths = partial_matches
    else:
        return f"No folder named or containing '{project_name}' found under {projects_dir}."

    abs_projects = os.path.abspath(projects_dir)
    project_paths = [
        p for p in project_paths
        if os.path.abspath(p).startswith(abs_projects)
    ]
    if not project_paths:
        return "Access denied: resolved paths are outside PROJECTS_DIR."

    result = []
    total_matches = 0
    all_abs_paths = []

    for project_path in project_paths:
        folder_name = os.path.basename(project_path)
        file_list = []
        abs_paths = []
        matches = []

        for root, dirs, files in os.walk(project_path):
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for fname in files:
                fpath = os.path.join(root, fname)
                rel_path = os.path.relpath(fpath, project_path)
                file_list.append((rel_path, fpath))
                abs_paths.append(fpath)
                if total_matches < 20:
                    try:
                        with open(fpath, encoding="utf-8", errors="ignore") as f:
                            for lineno, line in enumerate(f, 1):
                                if query.lower() in line.lower():
                                    matches.append(f"{rel_path}:{lineno}: {line.rstrip()}")
                                    total_matches += 1
                                    if total_matches >= 20:
                                        break
                    except OSError as e:
                        logger.debug("search_project_files: skipping %s: %s", fpath, e)

        all_abs_paths.extend(abs_paths)
        result.append(f"=== {folder_name} ({project_path}) ===")
        result.append(f"Files ({len(file_list)} total):")
        for rel, fpath in file_list[:50]:
            result.append(f"  {rel} ({_file_size_str(fpath)})")
        if len(file_list) > 50:
            result.append(f"  ... and {len(file_list) - 50} more files")
        if matches:
            result.append(f"Matches for '{query}':")
            result.extend(f"  {m}" for m in matches)
        else:
            result.append(f"No matches for '{query}'.")
        result.append("")

    if total_matches >= 20:
        result.append("(Results truncated at 20 total matches)")

    return "\n".join(result) + _build_file_select_block(all_abs_paths)


@tool
def find_files_by_name_exact(filename: str) -> str:
    """Find files whose name exactly matches the given filename (case-insensitive) across PROJECTS_DIR.
    Use this when the user knows the exact file name and wants to locate it on disk.
    filename: exact file name to search for, e.g. 'report.pdf', 'config.py'
    Returns a list of full file paths where the exact match was found.
    """
    projects_dir = os.getenv("PROJECTS_DIR")
    if not projects_dir:
        return "PROJECTS_DIR environment variable is not set."

    found = []
    for root, dirs, files in os.walk(projects_dir):
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fname in files + dirs:
            if fname.lower() == filename.lower():
                found.append(os.path.join(root, fname))
                if len(found) >= 20:
                    break
        if len(found) >= 20:
            break

    if not found:
        return f"No file named '{filename}' found under {projects_dir}."

    result = [f"Exact matches for '{filename}' ({len(found)}):"]
    for f in found:
        result.append(f"  {f} ({_file_size_str(f)})")
    if len(found) >= 20:
        result.append("  ... (truncated at 20 results)")
    return "\n".join(result) + _build_file_select_block(found)


@tool
def find_files_by_name_contains(keyword: str) -> str:
    """Find files whose name contains the given keyword (case-insensitive) across PROJECTS_DIR.
    Use this when the user wants to list all files that have a word or phrase in their filename.
    keyword: word or phrase to look for inside file names, e.g. 'blockchain', 'report', '2024'
    Returns a list of full file paths whose names contain the keyword.
    """
    projects_dir = os.getenv("PROJECTS_DIR")
    if not projects_dir:
        return "PROJECTS_DIR environment variable is not set."

    found = []
    for root, dirs, files in os.walk(projects_dir):
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fname in files + dirs:
            if keyword.lower() in fname.lower():
                found.append(os.path.join(root, fname))
                if len(found) >= 30:
                    break
        if len(found) >= 30:
            break

    if not found:
        return f"No files with '{keyword}' in their name found under {projects_dir}."

    result = [f"Files with '{keyword}' in name ({len(found)}):"]
    for f in found:
        result.append(f"  {f} ({_file_size_str(f)})")
    if len(found) >= 30:
        result.append("  ... (truncated at 30 results)")
    return "\n".join(result) + _build_file_select_block(found)


_dir_count_cache: dict[str, tuple[float, int]] = {}
_DIR_COUNT_CACHE_TTL = 300  # seconds


def _get_total_entries(projects_dir: str) -> int:
    """Total files+dirs under projects_dir, cached for _DIR_COUNT_CACHE_TTL seconds."""
    now = time.time()
    cached = _dir_count_cache.get(projects_dir)
    if cached and now - cached[0] < _DIR_COUNT_CACHE_TTL:
        return cached[1]

    total = 0
    for root, dirs, files in os.walk(projects_dir):
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        total += len(files) + len(dirs)
    _dir_count_cache[projects_dir] = (now, total)
    return total


def search_files_with_progress(
    projects_dir: str,
    filename: str | None,
    keyword: str | None,
    limit_exact: int = 20,
    limit_contains: int = 30,
):
    """Generator: walks projects_dir once looking for exact/contains name matches.
    Yields ("counting", None) if the entry-count cache needs refreshing, then
    ("progress", pct) as entries are scanned, then a final
    ("done", {"exact": [...], "contains": [...]})."""
    now = time.time()
    cached = _dir_count_cache.get(projects_dir)
    if not cached or now - cached[0] >= _DIR_COUNT_CACHE_TTL:
        yield ("counting", None)
    total_entries = _get_total_entries(projects_dir)

    visited = 0
    last_pct = -1
    exact_found = []
    contains_found = []

    for root, dirs, files in os.walk(projects_dir):
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fname in files + dirs:
            visited += 1

            if filename is not None and len(exact_found) < limit_exact and fname.lower() == filename.lower():
                exact_found.append(os.path.join(root, fname))
            if keyword is not None and len(contains_found) < limit_contains and keyword.lower() in fname.lower():
                contains_found.append(os.path.join(root, fname))

            pct = min(99, int(visited / total_entries * 100)) if total_entries else 99
            if pct > last_pct:
                last_pct = pct
                yield ("progress", pct)

        exact_done = filename is None or len(exact_found) >= limit_exact
        contains_done = keyword is None or len(contains_found) >= limit_contains
        if exact_done and contains_done:
            break

    yield ("done", {"exact": exact_found, "contains": contains_found})


def make_document_search_tool(retriever, embedding_model=None, sessions_dir: str = None):
    @tool
    def search_documents(query: str, config: RunnableConfig) -> str:
        """Search uploaded company documents for information, policies, procedures, and guidelines.
        Use this for ANY question about content in uploaded documents — product details, pricing,
        procedures, company policies, project information, or any other document content.
        Always search documents before answering questions about company-specific information.
        """
        docs = retriever.invoke(query)
        if sessions_dir and embedding_model:
            session_id = (config.get("configurable") or {}).get("thread_id")
            if session_id:
                session_dir = os.path.join(sessions_dir, session_id)
                if os.path.exists(session_dir):
                    session_vs = Chroma(persist_directory=session_dir, embedding_function=embedding_model)
                    session_docs = session_vs.as_retriever(search_kwargs={"k": 2}).invoke(query)
                    docs = docs + session_docs
        return "\n\n".join(d.page_content for d in docs) if docs else "No relevant information found in uploaded documents."
    return search_documents


class Tools:
    form_tools = [search_project_files, find_files_by_name_exact, find_files_by_name_contains]
    tools = [web_search, generate_presentation, generate_word_document, generate_pdf_document]
